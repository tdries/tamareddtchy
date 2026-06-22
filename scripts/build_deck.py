#!/usr/bin/env python3
"""Build the Tamareddtchy pitch deck. No master template was supplied by the
hackathon, so we build from a blank 16:9 deck in the incubator-toy brand.
Every run colors/sizes its own runs explicitly so nothing renders invisible.
Run: .venv/bin/python scripts/build_deck.py"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs", "screenshots")
OUT = os.path.join(ROOT, "docs", "Tamareddtchy-deck.pptx")

# brand palette
BG = RGBColor(0x15, 0x12, 0x1C)
PANEL = RGBColor(0x24, 0x1D, 0x31)
INK = RGBColor(0xF4, 0xEE, 0xFF)
INK2 = RGBColor(0xB9, 0xAE, 0xD0)
INK3 = RGBColor(0x84, 0x78, 0xA0)
ORANGE = RGBColor(0xFF, 0x45, 0x00)
GREEN = RGBColor(0x8E, 0xF0, 0xBD)
GLOW = RGBColor(0x7A, 0x5C, 0xFF)

DISPLAY = "Trebuchet MS"  # rounded, friendly; Fredoka not guaranteed in soffice
BODY = "Trebuchet MS"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs: list of paragraphs, each a list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = DISPLAY if bold else BODY
    return tb


def rect(s, x, y, w, h, color, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def eyebrow(s, x, y, label):
    text(s, x, y, Inches(6), Inches(0.4), [[(label.upper(), 13, ORANGE, True)]])


def bar(s, x, y):
    """the orange accent bar under titles"""
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.6), Pt(5))
    b.fill.solid()
    b.fill.fore_color.rgb = ORANGE
    b.line.fill.background()
    b.shadow.inherit = False


# ---------------------------------------------------------------- slide 1: cover
s = slide()
# glow disc
disc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.4), Inches(0.5), Inches(6.5), Inches(6.5))
disc.fill.solid(); disc.fill.fore_color.rgb = RGBColor(0x22, 0x1B, 0x33)
disc.line.fill.background(); disc.shadow.inherit = False
text(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.4),
     [[("Tamareddtchy", 64, INK, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
     [[("Your Reddit soul, raised as a living 3D pet.", 24, INK2, False)]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
     [[("Reddit's Games with a Hook  .  Devvit Web", 16, INK3, False)]], align=PP_ALIGN.CENTER)
bb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.87), Inches(5.4), Inches(1.6), Pt(6))
bb.fill.solid(); bb.fill.fore_color.rgb = ORANGE; bb.line.fill.background(); bb.shadow.inherit = False

# ---------------------------------------------------------------- slide 2: problem
s = slide()
eyebrow(s, Inches(0.9), Inches(0.7), "The problem")
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.0),
     [[("Reddit games die on day three.", 40, INK, True)]])
bar(s, Inches(0.95), Inches(2.05))
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2),
     [[("Nothing pulls you back. Nothing connects you to anyone else. You lurk the same "
        "three subreddits and a game post gets played once and scrolled past.", 20, INK2, False)]],
     spacing=1.15)
# two contrast cards
rect(s, Inches(0.9), Inches(3.9), Inches(5.6), Inches(2.7), PANEL, line=INK3)
rect(s, Inches(6.85), Inches(3.9), Inches(5.6), Inches(2.7), PANEL, line=ORANGE)
text(s, Inches(1.2), Inches(4.1), Inches(5.0), Inches(2.3),
     [[("Today", 18, INK3, True)],
      [("Engagement is an abstract number.", 16, INK2, False)],
      [("Winning is a score nobody remembers.", 16, INK2, False)],
      [("You never leave your bubble.", 16, INK2, False)]], spacing=1.3)
text(s, Inches(7.15), Inches(4.1), Inches(5.0), Inches(2.3),
     [[("With Tamareddtchy", 18, ORANGE, True)],
      [("Engagement is food that shapes a pet you love.", 16, INK, False)],
      [("Winning is a mutated lineage you built with others.", 16, INK, False)],
      [("Your best move is to find your opposite.", 16, INK, False)]], spacing=1.3)

# ---------------------------------------------------------------- slide 3: what it does
s = slide()
eyebrow(s, Inches(0.9), Inches(0.7), "What it does")
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.0),
     [[("A pet whose body IS your Reddit personality.", 34, INK, True)]])
bar(s, Inches(0.95), Inches(2.0))
pts = [
    ("Hatch", "Pick the corners of Reddit you live in. Your creature is born looking like you."),
    ("Feed", "Read, comment, post. Every action pushes points into a gene and reshapes the 3D body."),
    ("Grow", "It breathes, blinks, and droops when neglected. Egg to adult across days of nurture."),
    ("Breed", "Find your opposite and mate. Good genetics come from crossing distant gene pools."),
]
y = Inches(2.5)
for i, (h, b) in enumerate(pts):
    cx = Inches(0.9) + Inches(3.1) * i
    rect(s, cx, y, Inches(2.85), Inches(3.4), PANEL, line=INK3)
    text(s, cx + Inches(0.25), y + Inches(0.25), Inches(2.4), Inches(0.6),
         [[(h, 22, GREEN if i % 2 else GLOW, True)]])
    text(s, cx + Inches(0.25), y + Inches(1.0), Inches(2.4), Inches(2.2),
         [[(b, 15, INK2, False)]], spacing=1.2)

# ---------------------------------------------------------------- slide 4: the genome
s = slide()
eyebrow(s, Inches(0.9), Inches(0.7), "The mechanic")
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.0),
     [[("12 genes. 6 opposite pairs. One body, grown live.", 32, INK, True)]])
bar(s, Inches(0.95), Inches(2.0))
text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.9),
     [[("The creature is never stored as a picture. It is rebuilt in 3D every frame from "
        "twelve numbers. Each body slot maps to one opposite pair.", 19, INK2, False)]], spacing=1.15)
pairs = ["Knowledge / Vitality  .  head", "Tech / Heart  .  eyes", "Craft / Mayhem  .  torso",
         "Pulse / Lore  .  arms", "Inner / Social  .  legs", "Earth / Fiction  .  aura"]
for i, p in enumerate(pairs):
    col = i % 2
    row = i // 2
    cx = Inches(0.9) + Inches(6.0) * col
    cy = Inches(3.6) + Inches(1.0) * row
    rect(s, cx, cy, Inches(5.7), Inches(0.8), PANEL, line=INK3)
    text(s, cx + Inches(0.3), cy + Inches(0.12), Inches(5.2), Inches(0.6),
         [[(p, 17, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- slide 5: the economy
s = slide()
eyebrow(s, Inches(0.9), Inches(0.7), "The hook: a breeding economy")
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.0),
     [[("Opposites win. Twins make a busted, viral mess.", 32, INK, True)]])
bar(s, Inches(0.95), Inches(2.0))
rows = [
    ("Gen-1 x Gen-1", "a plain Gen-2", "cheap starter, breed wide"),
    ("Gen-3 x Gen-3", "a strong Gen-4", "both want the kid, hard-fought deal"),
    ("Gen-5 x Gen-1", "an advanced Gen-6", "the Gen-5 holder has the leverage"),
]
# header
rect(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.7), PANEL, line=ORANGE)
text(s, Inches(1.1), Inches(2.6), Inches(11), Inches(0.5),
     [[("Pairing            produces            and the strategy is", 16, ORANGE, True)]],
     anchor=MSO_ANCHOR.MIDDLE)
for i, (a, b, c) in enumerate(rows):
    cy = Inches(3.3) + Inches(0.85) * i
    rect(s, Inches(0.9), cy, Inches(11.5), Inches(0.72), PANEL, line=INK3)
    text(s, Inches(1.1), cy + Inches(0.1), Inches(3), Inches(0.5), [[(a, 16, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.3), cy + Inches(0.1), Inches(3), Inches(0.5), [[(b, 16, GREEN, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.6), cy + Inches(0.1), Inches(4.6), Inches(0.5), [[(c, 15, INK2, False)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.8),
     [[("Cost to mate: a big hunger hit plus a multi-day cooldown. Win metric: lineage score "
        "= depth x genetics quality + offspring count.", 16, INK3, False)]], spacing=1.15)

# ---------------------------------------------------------------- slide 6: the build (screenshots)
s = slide()
eyebrow(s, Inches(0.9), Inches(0.55), "See it")
text(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.8),
     [[("Every genome grows a different body.", 30, INK, True)]])
# nursery + mate side by side, framed
for path, x in [("01-nursery.png", Inches(0.9)), ("03-mate.png", Inches(6.95))]:
    p = os.path.join(SHOTS, path)
    if os.path.exists(p):
        frame = rect(s, x - Emu(12700), Inches(2.0) - Emu(12700), Inches(5.5) + Emu(25400), Inches(3.78) + Emu(25400), PANEL, line=INK3)
        s.shapes.add_picture(p, x, Inches(2.0), width=Inches(5.5))
text(s, Inches(0.9), Inches(6.0), Inches(5.5), Inches(0.6), [[("The nursery: a living 3D pet.", 15, INK2, False)]], align=PP_ALIGN.CENTER)
text(s, Inches(6.95), Inches(6.0), Inches(5.5), Inches(0.6), [[("The mate market: five genomes, five bodies.", 15, INK2, False)]], align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- slide 7: tech
s = slide()
eyebrow(s, Inches(0.9), Inches(0.7), "How it is built")
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.0),
     [[("One Devvit Web app. Shared rules. No model files.", 30, INK, True)]])
bar(s, Inches(0.95), Inches(1.95))
techs = [
    ("Procedural 3D", "Three.js geometry grown from the genome. No assets to download, infinite variety."),
    ("Shared logic", "Genome, scoring, and breeding live in one module imported by client, server, and tests."),
    ("Devvit Web", "A web view post backed by Redis and the Reddit API. Milestones post themselves."),
    ("Plays anywhere", "An in-memory mock makes the whole game playable with zero Reddit auth, for judges."),
]
for i, (h, b) in enumerate(techs):
    col = i % 2; row = i // 2
    cx = Inches(0.9) + Inches(6.0) * col
    cy = Inches(2.5) + Inches(1.95) * row
    rect(s, cx, cy, Inches(5.7), Inches(1.7), PANEL, line=INK3)
    text(s, cx + Inches(0.3), cy + Inches(0.2), Inches(5.1), Inches(0.5), [[(h, 19, GLOW, True)]])
    text(s, cx + Inches(0.3), cy + Inches(0.75), Inches(5.1), Inches(0.9), [[(b, 15, INK2, False)]], spacing=1.15)

# ---------------------------------------------------------------- slide 8: closing
s = slide()
disc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.4), Inches(0.5), Inches(6.5), Inches(6.5))
disc.fill.solid(); disc.fill.fore_color.rgb = RGBColor(0x22, 0x1B, 0x33)
disc.line.fill.background(); disc.shadow.inherit = False
text(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
     [[("Leave your bubble.", 52, INK, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.7),
     [[("Raise something only the whole of Reddit could make.", 22, INK2, False)]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(4.7), Inches(11.3), Inches(0.6),
     [[("github.com/tdries/tamareddtchy", 17, ORANGE, True)]], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
